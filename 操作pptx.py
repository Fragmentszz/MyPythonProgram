#encoding=gbk

from pptx import Presentation
import requests
from excel import zqsport,ydxx,count
import re
from pptx.dml.color import RGBColor
from pptx.util import Pt,Cm


def min2hmin(minz):
    h=minz//60
    minz=minz % 60
    if h==0:
        return '{}min'.format(minz)
    elif minz==0:
        return '{}h'.format(h)
    else:
        return '{}h{}min'.format(h,minz)




def mc(list:list):
    dic={"0":"一","1":"二","2":"三","3":"四","4":"五"}
    count=0
    result={"第一名":[0,''],"第二名":[0,''],"第三名":[0,''],"第四名":[0,''],"第五名":[0,'']}
    flag=list[0][1]
    i=0
    mc="第{}名".format(dic[str(i)])
    result[mc][0]=flag
    for tp in list:
        if tp[1]==flag:
            pass
        else:
            i+=1
            if i==5:
                break
            mc="第{}名".format(dic[str(i)])
            result[mc][0]=tp[1]
        result[mc][1]=result[mc][1]+' {}'.format(tp[0])
        flag=tp[1]
    
    return result

def rq(rq:str):
    list=[]
    rq_list=rq.split('-')
    for rz in rq_list:
        p1='(\d+).'
        p2='\d.(\d+)'
        r1=eval(re.findall(p1,rz)[0].strip('0'))
        r2=eval(re.findall(p2,rz)[0].strip('0'))
        list.append(r1)
        list.append(r2)
    return list


def mcc (time0,sort:str):
    li=zqsport(time0,sort)
    dic=mc(li)
    dic2={"0":"一","1":"二","2":"三","3":"四","4":"五"}
    i=0
    zz="第{}名".format(dic2[str(i)])
    str_list=[]
    while dic[zz][0]!=0:
        str_list.append("{}（打卡{}次）：{}".format(zz,dic[zz][0],dic[zz][1]))
        i+=1
        try:
            zz="第{}名".format(dic2[str(i)])
        except:
            break
    return str_list

def mcz(time0,sort:str):
    li=ydxx(time0,sort)
    dic=mc(li)
    dic2=count(time0,sort)
    dic3={"0":"一","1":"二","2":"三","3":"四","4":"五"}
    i=0
    zz="第{}名".format(dic3[str(i)])
    str_list=[]
    while dic[zz][0]!=0:
        name_list=dic[zz][1].split()
        strz=''
        if sort=="学习":
            for name in name_list:
                strz=strz+'           {}(学习总时长{}，打卡{}次)'.format(name,min2hmin(dic[zz][0]),dic2[name])+'\n'
        if sort=="阅读":
            for name in name_list:
                strz=strz+'           (阅读总时长{})：{}(打卡{}次)'.format(min2hmin(dic[zz][0]),name,dic2[name])+'\n'
        strz='{}:'.format(zz)+strz.strip(' ').strip('\n')
        str_list.append(strz)
        i+=1
        try:
            zz="第{}名".format(dic3[str(i)])
        except:
            break
    return str_list

def xieru(shape,name,bold,italic,rgb,size,str_list):
    z=0
    print(shape.text_frame.text)
    count=len(shape.text_frame.paragraphs)
    for strr in str_list:
        try:
            paragraph=shape.text_frame.paragraphs[z]
        except Exception as e:
            paragraph=shape.text_frame.add_paragraph()
            print(e)
        paragraph.text=strr
        font=paragraph.font
        font.name=name
        font.bold=bold
        font.italic=italic
        font.color.rgb=RGBColor(*rgb)
        font.size=size
        z+=1
    if count>z:
        for i in range(z,count):
            paragraph=shape.text_frame.paragraphs[i]
            paragraph.text=''

def font(shape,flag,sort):
    tf=shape.text_frame
    alignment=tf.paragraphs[0].alignment # 对齐方式

    bold=False   # 是否加粗
    italic=False   # 是否斜体
    if flag:
        rgb=(0,0,0)  # 字体颜色
        if sort=="阅读" or sort=="学习":
            size=Pt(20)
        else:
            size=Pt(24)
        name="汉仪天宇风行体W"  # 字体名称
    else:
        rgb=(254,233,177)
        size=Pt(28)                  # 字体大小
        name="华文新魏"  # 字体名称
    return (alignment,bold,italic,rgb,size,name)


def cz(shapes_list,time0,sort):
    for j in range(4):
        shape=shapes_list[j]
        if shape.has_text_frame and shape.text[-2:]!="打卡":

            if shape.text.strip('\n')[:3]=="“悦读":
                alignment,bold,italic,rgb,size,name=font(shape,0,sort)
                li=rq(time0)
                str_list=["“悦读·跃动”寒假线上打卡排行榜{}月{}日至{}月{}日".format(li[0],li[1],li[2],li[3])]
            else:
                alignment,bold,italic,rgb,size,name=font(shape,1,sort)
                if sort=="早起" or sort=="运动":
                    str_list=mcc(time0,sort)
                if sort=="阅读" or sort=="学习":
                    str_list=mcz(time0,sort)
            xieru(shape,name,bold,italic,rgb,size,str_list)


    

            

    





def pptx(time0):
    prs=Presentation('ss.pptx')
    for i in range(4):
        slide=prs.slides[i]
        shapes_list=list(slide.shapes)
        
        if i ==0:
            cz(shapes_list,time0,"早起")
        if i==1:
            cz(shapes_list,time0,"运动")
        if i==2:
            cz(shapes_list,time0,"学习")
        if i==3:
            cz(shapes_list,time0,"阅读")
                                
    prs.save('ss.pptx')
pptx('01.08-02.20')

            
      

